"""Generate PowerPoint slides from parsed manual content."""

from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DEFAULT_COLOR_SCHEME: dict[str, tuple[int, int, int]] = {
    "primary": (0, 102, 204),
    "accent": (51, 153, 255),
    "dark": (0, 51, 102),
    "light_bg": (240, 248, 255),
    "gray": (100, 100, 100),
}


def _to_rgb(rgb: tuple[int, int, int]) -> RGBColor:
    return RGBColor(rgb[0], rgb[1], rgb[2])


class PowerPointGenerator:
    """Generate a CellGenTech-styled PowerPoint presentation."""

    def __init__(
        self,
        output_path: str,
        logo_path: str | None = None,
        color_scheme: dict[str, tuple[int, int, int]] | None = None,
    ) -> None:
        self.output_path = output_path
        self.logo_path = logo_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        scheme = dict(DEFAULT_COLOR_SCHEME)
        if color_scheme:
            for key, value in color_scheme.items():
                if key in scheme:
                    scheme[key] = value

        self.color_primary = _to_rgb(scheme["primary"])
        self.color_accent = _to_rgb(scheme["accent"])
        self.color_dark = _to_rgb(scheme["dark"])
        self.color_light_bg = _to_rgb(scheme["light_bg"])
        self.color_gray = _to_rgb(scheme["gray"])

    def add_logo_to_slide(self, slide: Any) -> None:
        """Add logo to the bottom-left corner if available."""
        if self.logo_path and os.path.exists(self.logo_path):
            try:
                logo_height = Inches(0.4)
                logo_left = Inches(0.3)
                logo_top = self.prs.slide_height - logo_height - Inches(0.2)
                slide.shapes.add_picture(self.logo_path, logo_left, logo_top, height=logo_height)
            except Exception as exc:  # noqa: BLE001
                print(f"WARN: failed to add logo: {exc}")

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), self.prs.slide_width, Inches(2.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.color_primary
        bg_shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = "Meiryo"
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.name = "Meiryo"
            subtitle_para.font.color.rgb = self.color_dark
            subtitle_para.alignment = PP_ALIGN.CENTER

        accent_line = slide.shapes.add_shape(1, Inches(2), Inches(4.0), Inches(6), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.color_accent
        accent_line.line.fill.background()

        self.add_logo_to_slide(slide)

    def add_content_slide(
        self,
        step_number: int,
        step_text: str,
        image_path: str | None,
        additional_images: list[dict[str, Any]] | None = None,
    ) -> None:
        """Add one content slide for a manual step."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        header_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), self.prs.slide_width, Inches(0.6))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.color_primary
        header_bg.line.fill.background()

        step_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(3), Inches(0.4))
        text_frame = step_box.text_frame
        text_frame.text = f"STEP {step_number}"
        para = text_frame.paragraphs[0]
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.name = "Meiryo"
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.LEFT

        desc_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(0.7), Inches(9.2), Inches(0.8))
        desc_bg.fill.solid()
        desc_bg.fill.fore_color.rgb = self.color_light_bg
        desc_bg.line.color.rgb = self.color_accent
        desc_bg.line.width = Pt(1)

        desc_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(8.8), Inches(0.5))
        desc_frame = desc_box.text_frame
        desc_frame.text = step_text
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(20)
        desc_para.font.name = "Meiryo"
        desc_para.font.color.rgb = self.color_dark
        desc_para.alignment = PP_ALIGN.LEFT

        if image_path and os.path.exists(image_path):
            img_top = Inches(1.7)
            img_left = Inches(0.5)
            max_width = Inches(9)
            max_height = Inches(5.2)
            try:
                pic = slide.shapes.add_picture(image_path, img_left, img_top, width=max_width)
                if pic.height > max_height:
                    aspect_ratio = pic.width / pic.height
                    pic.height = max_height
                    pic.width = int(max_height * aspect_ratio)
                pic.left = int((self.prs.slide_width - pic.width) / 2)
            except Exception as exc:  # noqa: BLE001
                print(f"WARN: failed to add image: {exc}")

        if additional_images:
            # Additional images are currently acknowledged but not rendered
            # to preserve the existing behavior.
            _ = additional_images

        self.add_logo_to_slide(slide)

    def save(self) -> bool:
        """Save presentation."""
        try:
            self.prs.save(self.output_path)
            return True
        except Exception as exc:  # noqa: BLE001
            print(f"ERROR: failed to save PowerPoint: {exc}")
            return False


def main() -> int:
    """Legacy standalone generator entrypoint."""
    module_path = Path(__file__).resolve()
    repo_root = module_path.parents[3]
    output_path = repo_root / "05_mail" / "相見積操作マニュアル.pptx"

    generator = PowerPointGenerator(str(output_path))
    generator.add_title_slide(title="相見積操作マニュアル", subtitle="見積依頼の作成と送信")
    generator.add_content_slide(step_number=1, step_text="サンプルステップ", image_path=None)
    return 0 if generator.save() else 1


if __name__ == "__main__":
    raise SystemExit(main())

