"""
PowerPoint生成スクリプト

このスクリプトは、Excelマニュアルから抽出した情報をもとにPowerPointスライドを生成します。
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Windows環境でのUnicodeエンコーディング設定
if sys.platform == 'win32':
    import codecs
    if hasattr(sys.stdout, 'buffer'):
        sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
        sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')


class PowerPointGenerator:
    """PowerPointファイルを生成するクラス"""

    def __init__(self, output_path, logo_path=None):
        """
        初期化

        Args:
            output_path: 出力PowerPointファイルのパス
            logo_path: ロゴ画像のパス（オプション）
        """
        self.output_path = output_path
        self.logo_path = logo_path
        self.prs = Presentation()
        # スライドサイズ: 標準16:9
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # カラースキーム（CellGenTechブルーに合わせる）
        self.color_primary = RGBColor(0, 102, 204)  # CellGenTechブルー
        self.color_accent = RGBColor(51, 153, 255)  # ライトブルー
        self.color_dark = RGBColor(0, 51, 102)  # ダークブルー
        self.color_light_bg = RGBColor(240, 248, 255)  # 薄い青背景
        self.color_gray = RGBColor(100, 100, 100)  # グレー

    def add_logo_to_slide(self, slide):
        """
        スライドの左下にロゴを追加

        Args:
            slide: ロゴを追加するスライド
        """
        if self.logo_path and os.path.exists(self.logo_path):
            try:
                # ロゴのサイズと位置
                logo_height = Inches(0.4)  # ロゴの高さ
                logo_left = Inches(0.3)  # 左からの位置
                logo_top = self.prs.slide_height - logo_height - Inches(0.2)  # 下からの位置

                # ロゴを追加
                pic = slide.shapes.add_picture(
                    self.logo_path,
                    logo_left,
                    logo_top,
                    height=logo_height
                )
            except Exception as e:
                print(f"    ! ロゴの追加に失敗: {e}")

    def add_title_slide(self, title, subtitle=""):
        """
        タイトルスライドを追加

        Args:
            title: メインタイトル
            subtitle: サブタイトル
        """
        print(f"  タイトルスライドを作成中...")

        # 空白スライドを使用してカスタムデザイン
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 背景に青のグラデーション風の四角形を追加
        bg_shape = slide.shapes.add_shape(
            1,  # 四角形
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(2.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.color_primary
        bg_shape.line.fill.background()

        # タイトル（白文字で青背景の上に）
        title_left = Inches(0.5)
        title_top = Inches(0.8)
        title_width = Inches(9)
        title_height = Inches(1.2)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.name = 'メイリオ'
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # 白
        title_para.alignment = PP_ALIGN.CENTER

        # サブタイトル（青背景の下に）
        if subtitle:
            subtitle_left = Inches(0.5)
            subtitle_top = Inches(3.0)
            subtitle_width = Inches(9)
            subtitle_height = Inches(0.8)
            subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.name = 'メイリオ'
            subtitle_para.font.color.rgb = self.color_dark
            subtitle_para.alignment = PP_ALIGN.CENTER

        # アクセントライン（装飾）
        accent_line = slide.shapes.add_shape(
            1,  # 四角形
            Inches(2), Inches(4.0),
            Inches(6), Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.color_accent
        accent_line.line.fill.background()

        # ロゴを追加
        self.add_logo_to_slide(slide)

        print(f"    ✓ タイトルスライド作成完了")

    def add_content_slide(self, step_number, step_text, image_path, additional_images=None):
        """
        コンテンツスライドを追加

        Args:
            step_number: ステップ番号
            step_text: ステップの説明文
            image_path: メイン画像のパス
            additional_images: 追加画像のリスト（オプション）
        """
        print(f"  STEP {step_number} スライドを作成中...")

        # 空白スライドレイアウト（レイアウト6）
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # ヘッダー背景（青のバー）
        header_bg = slide.shapes.add_shape(
            1,  # 四角形
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.6)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.color_primary
        header_bg.line.fill.background()

        # ステップ番号（白文字で青背景の上に）
        step_left = Inches(0.5)
        step_top = Inches(0.1)
        step_width = Inches(3)
        step_height = Inches(0.4)
        step_box = slide.shapes.add_textbox(step_left, step_top, step_width, step_height)
        text_frame = step_box.text_frame
        text_frame.text = f"STEP {step_number}"
        para = text_frame.paragraphs[0]
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.name = 'メイリオ'
        para.font.color.rgb = RGBColor(255, 255, 255)  # 白
        para.alignment = PP_ALIGN.LEFT

        # 説明文を追加（背景の薄い青の枠）
        desc_bg_left = Inches(0.4)
        desc_bg_top = Inches(0.7)
        desc_bg_width = Inches(9.2)
        desc_bg_height = Inches(0.8)

        desc_bg = slide.shapes.add_shape(
            1,  # 四角形
            desc_bg_left, desc_bg_top,
            desc_bg_width, desc_bg_height
        )
        desc_bg.fill.solid()
        desc_bg.fill.fore_color.rgb = self.color_light_bg
        desc_bg.line.color.rgb = self.color_accent
        desc_bg.line.width = Pt(1)

        # 説明文テキスト
        desc_left = Inches(0.6)
        desc_top = Inches(0.85)
        desc_width = Inches(8.8)
        desc_height = Inches(0.5)
        desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
        desc_frame = desc_box.text_frame
        desc_frame.text = step_text
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(20)
        desc_para.font.name = 'メイリオ'
        desc_para.font.color.rgb = self.color_dark
        desc_para.alignment = PP_ALIGN.LEFT

        # メイン画像を追加
        if image_path and os.path.exists(image_path):
            img_top = Inches(1.7)
            img_left = Inches(0.5)
            max_width = Inches(9)
            max_height = Inches(5.2)

            try:
                # 画像を追加（アスペクト比を維持）
                pic = slide.shapes.add_picture(
                    image_path,
                    img_left,
                    img_top,
                    width=max_width
                )

                # 高さがmax_heightを超える場合は調整
                if pic.height > max_height:
                    # アスペクト比を維持して縮小
                    aspect_ratio = pic.width / pic.height
                    pic.height = max_height
                    pic.width = int(max_height * aspect_ratio)

                # 中央揃え
                pic.left = int((self.prs.slide_width - pic.width) / 2)

                print(f"    ✓ 画像を追加: {os.path.basename(image_path)}")
            except Exception as e:
                print(f"    ✗ 画像の追加に失敗: {e}")

        # 追加画像がある場合
        if additional_images:
            print(f"    追加画像: {len(additional_images)}件")
            # 追加画像は次のスライドに配置するか、小さく配置
            # ここでは簡略化のため省略

        # ロゴを追加
        self.add_logo_to_slide(slide)

        print(f"    ✓ STEP {step_number} スライド作成完了")

    def save(self):
        """PowerPointファイルを保存"""
        try:
            self.prs.save(self.output_path)
            print(f"\n✓ PowerPointファイルを保存しました: {self.output_path}")
            return True
        except Exception as e:
            print(f"\n✗ エラー: PowerPointファイルの保存に失敗しました: {e}")
            return False


def main():
    """メイン処理"""
    print("=" * 60)
    print("PowerPoint生成スクリプト")
    print("=" * 60)

    # 出力パス
    script_dir = Path(__file__).parent
    project_dir = script_dir.parent
    output_path = project_dir / "相見積操作マニュアル.pptx"

    print(f"\n出力ファイル: {output_path}")

    # PowerPointジェネレーターを初期化
    generator = PowerPointGenerator(str(output_path))

    # タイトルスライドを追加
    generator.add_title_slide(
        title="相見積操作マニュアル",
        subtitle="見積依頼の作成と送信"
    )

    # サンプルスライドを追加（実際のデータは統合スクリプトから渡される）
    print("\nサンプルスライドを作成中...")
    generator.add_content_slide(
        step_number=1,
        step_text="サンプルステップです",
        image_path=None
    )

    # 保存
    print("\n" + "=" * 60)
    print("PowerPointファイルを保存中...")
    print("=" * 60)
    if generator.save():
        print(f"\n✓ PowerPoint生成が正常に完了しました")
        return 0
    else:
        print(f"\n✗ PowerPoint生成に失敗しました")
        return 1


if __name__ == "__main__":
    sys.exit(main())
